import shutil

import openai
import uuid
import firebase_admin
import os
import base64
import bcrypt
import json
import io

from flask import Flask, request, redirect, url_for, flash, session, send_from_directory, send_file, jsonify
from views import *
from firebase_admin import credentials, firestore, storage
from flask_login import login_required, logout_user
from utils.generate import parse_response, create_ppt, update_slide_ppt
from utils.chatdev import chat_development, slide_chat_development
from pptx import Presentation
import aspose.slides as slides
import aspose.pydrawing as drawing
from spire.presentation import *
from pymongo.mongo_client import MongoClient
from pymongo.server_api import ServerApi
from io import BytesIO
from bson import ObjectId, json_util

app: Flask = Flask(__name__)
# app.config['PERMANENT_SESSION_LIFETIME'] = 3600
app.register_blueprint(home)
app.register_blueprint(landing)
app.register_blueprint(account)
app.register_blueprint(choose)
app.register_blueprint(generate)
app.register_blueprint(presentation)
app.register_blueprint(test)
app.register_blueprint(choosetemplate)

# Generate a random secret key
app.secret_key = os.urandom(24)

# Your web app's Firebase configuration
firebase_config = {
    "apiKey": "AIzaSyBl2jxG_tenXKAX2NySytDXgw3PSoqaci0",
    "authDomain": "smartsync-ade70.firebaseapp.com",
    "projectId": "smartsync-ade70",
    "storageBucket": "smartsync-ade70.appspot.com",
    "messagingSenderId": "602190235087",
    "appId": "1:602190235087:web:dcc245609c0ccd2f1d3eba",
    "measurementId": "G-MWF9ER8V2Z"
}

# OpenAI API key
openai.api_key = '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'

# # Your Firebase configuration
# cred = credentials.Certificate("D:\RitchelMendaros\PyCharm_Projects\smartsync-ade70-firebase-adminsdk-l2ti0-1ea8a94791.json")
# firebase_admin.initialize_app(cred)
#
# # Initialize Firebase Storage
# firebase_storage = storage.bucket(app=firebase_admin.get_app(), name="smartsync-ade70.appspot.com")
#
# # Reference to the Firestore database
# db = firestore.client()
#
# os.makedirs('generated', exist_ok=True)

# Connect to MONGODB
uri = "mongodb+srv://thertj:amara@smartsync.kafshog.mongodb.net/?retryWrites=true&w=majority"

# Create a new client and connect to the server
client = MongoClient(uri, server_api=ServerApi('1'))

db = client['test2']

try:
    client.admin.command('ping')
    print("Pinged your deployment. You successfully connected to MongoDB!")
except Exception as e:
    print(e)


@app.route('/')
def default():
    return render_template('LandingPage.html')


collection = 'user_collections'
user_collections = db[collection]
# Sign-in route
@app.route('/signin', methods=['POST'])
def signin():
    if request.method == 'POST':
        username = request.form.get('username')
        password = request.form.get('password')

        # Check if the username exists in the user collection
        user = user_collections.find_one({'username': username})
        try:
            if user:
                # Check if the password is correct
                if bcrypt.checkpw(password.encode('utf-8'), user['password']):
                    # Successfully signed in, set session variable or perform other actions
                    session['username'] = username
                    flash('Login successful. Welcome!', 'success')
                    return redirect(url_for('search_collections'))
                else:
                    flash('Invalid username or password', 'error')
                    return render_template('SignIn_UI.html')
            else:
                flash('User not found!', 'error')
                return render_template('SignIn_UI.html')
        except Exception as e:
            # Error during authentication
            print(f"Error during sign-in: {e}")
            flash('Error during sign-in', 'error')
            return render_template('SignIn_UI.html')
    # Handle other cases (GET request, etc.)
    return redirect(url_for('default'))


# Sign-up route
@app.route('/signup', methods=['GET', 'POST'])
def signup():

    if request.method == 'POST':
        username = request.form.get('username')
        password = request.form.get('password')
        first_name = request.form.get('first_name')
        last_name = request.form.get('last_name')
        email = request.form.get('email')
        contact_number = request.form.get('contact_number')
        hashed_password = bcrypt.hashpw(password.encode('utf-8'), bcrypt.gensalt())

        # Check for non-empty fields
        if not all([username, password, first_name, last_name, email, contact_number]):
            flash('All fields must be filled out.', 'error')
            return render_template('SignUp_UI.html')

        # Check email format
        if not is_valid_email(email):
            flash('Invalid email address.', 'error')
            return render_template('SignUp_UI.html')

        # Check contact number format
        if not is_valid_contact_number(contact_number):
            flash('Invalid contact number.', 'error')
            return render_template('SignUp_UI.html')

        try:
            # Check if the username is already taken
            if user_collections.find_one({'username': username}):
                flash('Username already taken. Please choose another.', 'error')
                return redirect(url_for('signup'))

            # Create a new user document
            new_user = {
                'username': username,
                'password': hashed_password,
                'first_name': first_name,
                'last_name': last_name,
                'email': email,
                'contact_number': contact_number
            }

            # Insert the new user into the user collection
            user_collections.insert_one(new_user)

            # Successful sign-up, you can add session handling or redirect as needed
            flash('Signup successful! You can now log in.', 'success')
            session['username'] = username
            return redirect(url_for('search_collections'))
        except Exception as e:
            # Error during user creation
            print(f"Error during sign-up: {e}")
            flash('Error during sign-up', 'error')
            return render_template('SignUp_UI.html')
    return render_template('SignUp_UI.html')


# checks if inputted email has @
def is_valid_email(email):
    return '@' in email


# checks if inputted contact_number is a digit
def is_valid_contact_number(contact_number):
    return contact_number.isdigit()


# Context processor to make user_info(firstname) available to all templates
@app.context_processor
def inject_user_info():
    user_info = {'first_name': None, 'username': None, 'is_logged_in': False}
    # Check if the user is authenticated
    if 'username' in session:
        username = session['username']
        user_info['is_logged_in'] = True
        # Retrieve user information from the 'user_collections' collection using the username
        user = user_collections.find_one({'username': username})
        if user:
            user_info['first_name'] = user.get('first_name')
            user_info['username'] = username
    return {'user_info': user_info}


# Get response for GenerateKeyPoints
@app.route('/get_response', methods=['POST'])
def get_response():
    topic = request.form['topic']
    num_slides = request.form.get('num_slides')
    objectives = request.form['objectives']
    prompt = (
        f"You are an expert in {topic}. Write the slide titles for a powerpoint presentation covering the "
        f"following topics and objectives {objectives}. Make it {num_slides} slides."
    )
    # Generate a response from OpenAI's GPT-3
    response = openai.Completion.create(
        engine="text-davinci-002",
        prompt=prompt,
        max_tokens=1000  # You can adjust the token limit as needed
    )
    bot_response = response.choices[0].text.strip()
    # Split the response into slide contents
    slide_contents = bot_response.split('\n')
    # Initialize the structured response with the title slide
    structured_response = f"\nTopic: {topic}\n\n"
    # Add slide content for the subsequent slides
    current_slide = 0  # Start with the second slide
    for i in range(len(slide_contents)):
        content = slide_contents[i].strip()
        if content and content[0].isdigit():
            # New slide detected
            current_slide += 1
            structured_response += f"Slide {current_slide}: {content.split(maxsplit=1)[1]}\n"
        else:
            structured_response += f"{content}\n"
    return structured_response


@app.route('/ChooseTemplate')
def choose_template():
    # Add logic for rendering the ChooseTemplate page
    return render_template('ChooseTemplate.html')


@app.route('/GeneratePresentation', methods=['GET','POST'])
def generate_presentation():
    # if request.method == 'GET':
    template = session.get('selected_template')
    if template is None:
        # If template_choice is None, redirect to ChooseTemplate
        return redirect(url_for('choose_template'))
    if request.method == 'POST':
        action = request.form.get('action')
        if action == 'action1':
            username = session.get('username')
            content = request.form.get('contents')
            template_choice = session.get('selected_template')
            print(template_choice)
            topic = request.form.get('title')
            presentor = request.form.get('presentation-presentor')
            cleaned_topic = topic.strip()

            if template_choice == "simple":
                template_choice = "simple"
            elif template_choice == "dark_modern":
                template_choice = "dark_modern"
            elif template_choice == "minimal_darkgreen":
                template_choice = "minimal_darkgreen"
            elif template_choice == "minimal_blue":
                template_choice = "minimal_blue"
            elif template_choice == "minimal_neon":
                template_choice = "minimal_neon"
            elif template_choice == "minimal_gray":
                template_choice = "minimal_gray"

            assistant_response = chat_development(content)
            session['assistant_response'] = assistant_response
            session['template'] = template_choice
            session['topic'] = cleaned_topic
            print(assistant_response)
            slides_content = parse_response(assistant_response)
            print(slides_content)

            try:
                ppt_filename = create_ppt(slides_content, template_choice, topic, presentor)
                ppt_path = os.path.join(os.path.abspath('generated'), 'generated_presentation.pptx')

                # Convert the PowerPoint presentation to images
                images_folder = os.path.join('static', cleaned_topic)
                if os.path.exists(images_folder):
                    # If it exists, delete it
                    shutil.rmtree(images_folder)

                os.makedirs(images_folder)
                image_paths = convert_ppt_to_images(ppt_path, images_folder, cleaned_topic, username)
                image_files = slideshow(images_folder, cleaned_topic)
                return render_template('GeneratePresentation.html', image_files=image_files)
            except Exception as e:
                flash(f'Error: {e}', 'error')
                return render_template('GeneratePresentation.html')

        elif action == 'action2':
            # try:
            slideNum = request.form.get('slide_num')
            instruction = request.form.get('instruction')
            uploaded_file = request.files.get('filename')
            is_auto_generated = 'isAuto' in request.form
            retain = True
            print(uploaded_file)
            if uploaded_file:
                hasPicture = True
            else:
                hasPicture = False

            file_path = ""
            if uploaded_file and uploaded_file.filename:
                retain = False
                if 'filename' not in request.files:
                    return "No file part"

                uploaded_file = request.files['filename']
                # If the user does not select a file, the browser submits an empty file without a filename
                if uploaded_file.filename == '':
                    return "No selected file"
                # Save the file to a location (replace 'uploads' with your desired directory)
                upload_folder = 'uploads'
                if not os.path.exists(upload_folder):
                    os.makedirs(upload_folder)

                file_path = os.path.join(upload_folder, uploaded_file.filename)
                uploaded_file.save(file_path)

            slide_content_from_session = session.get('assistant_response')
            template_choice = session.get('template')
            topic = session.get('topic')

            assistant_response = slide_chat_development(slide_content_from_session, instruction, slideNum)
            print(assistant_response)
            slides_content = parse_response(assistant_response)
            print(slides_content)
            update_slide_ppt(slides_content, file_path, is_auto_generated, hasPicture, template_choice, slideNum, retain)
            update_slide_ppt(slides_content, file_path, is_auto_generated, hasPicture, template_choice, slideNum, retain)

            ppt_path = os.path.join(os.path.abspath('generated'), 'generated_presentation.pptx')
            folder_path = f"static/{topic}"
            convert_updated_slide(ppt_path, folder_path, topic, int(slideNum))
            images_folder = os.path.join('static', topic)
            image_files = slideshow(images_folder, topic)
            return render_template('GeneratePresentation.html', image_files=image_files, action='action1')


def slideshow(images_folder, topic):
    # Get the list of image files in the static folder
    image_files = [f"static/{topic}/{file}" for file in os.listdir(f"static/{topic}") if file.endswith('.png')]

    # Extract the numerical part of the filename and use it as the sorting key
    def extract_number(filename):
        return int(''.join(filter(str.isdigit, filename)))

    # Sort the images based on the numerical part of their filenames
    image_files.sort(key=extract_number)

    # Render the template with the list of image files
    return image_files

@app.route('/slideshow1', methods=['POST'])
def slideshow1(topic):
    images_folder = 'static'
    # Get the list of image files in the static folder
    image_files = [f"static/{topic}/{file}" for file in os.listdir(f"static/{topic}") if file.endswith('.png')]

    # Extract the numerical part of the filename and use it as the sorting key
    def extract_number(filename):
        return int(''.join(filter(str.isdigit, filename)))

    # Sort the images based on the numerical part of their filenames
    image_files.sort(key=extract_number)

    # Render the template with the list of image files
    return image_files

@app.route('/select_template', methods=['POST'])
def select_template():
    if request.method == 'POST':
        template_name = request.form.get('template_name')
        print(template_name)
        # Perform actions with the selected template, e.g., store it in the session
        session['selected_template'] = template_name
        return render_template('GeneratePresentation.html', template_name=template_name)


def convert_ppt_to_images(ppt_path, output_folder, topic, username):
    # Connect to MongoDB
    client = MongoClient("mongodb+srv://thertj:amara@smartsync.kafshog.mongodb.net/?retryWrites=true&w=majority")
    db = client['test2']
    new_collection = f'{username}_{topic}'

    if new_collection not in db.list_collection_names():
        db.create_collection(new_collection)

    image_collection = db[new_collection]

    presentation = Presentation()
    print(output_folder)

    try:
        # Load the PowerPoint presentation
        presentation.LoadFromFile(ppt_path)

        # Loop through the slides in the presentation
        image_paths = []
        for i, slide in enumerate(presentation.Slides):
            # Specify the output file name
            file_name = f"{output_folder}/{topic}_slide_{i + 1}.png"
            # Save each slide as a PNG image
            image = slide.SaveAsImage()
            image.Save(file_name)
            image_paths.append(file_name)
            image.Dispose()

            # Save the image path to MongoDB
            image_collection.insert_one({
                'topic': topic,
                'slide_number': i + 1,
                'image_path': file_name
            })
        print("success")
        return image_paths

    finally:
        # Dispose of the presentation object
        presentation.Dispose()
        # Close the MongoDB connection
        client.close()


# Insert image to MONGODB
def insert_image(image_path):
    with open(image_path, 'rb') as image_file:
        image_binary = image_file.read()
        document = {
            "filename": os.path.basename(image_path),
            "image_data": image_binary,
        }
        result = collection.insert_one(document)
        return result.inserted_id

def convert_updated_slide(ppt_path, output_folder, topic, slideNum):
    presentation = Presentation()
    print(output_folder)
    try:
        # Load the PowerPoint presentation
        presentation.LoadFromFile(ppt_path)
        i = 0
        # Loop through the slides in the presentation
        image_paths = []
        for slide in presentation.Slides:
            # Specify the output file name
            if i == slideNum:
                file_name = f"{output_folder}/{topic}_slide_{i + 1}.png"
                # Save each slide as a PNG image
                if os.path.exists(file_name):
                    # If the file exists, remove it before saving the updated image
                    os.remove(file_name)
                image = slide.SaveAsImage()
                image.Save(file_name)
                image_paths.append(file_name)
                image.Dispose()
                return image_paths
            i += 1
        print("success")
        return image_paths

    finally:
        # Dispose of the presentation object
        presentation.Dispose()


#Custom JSON Encoder
class MongoJSONEncoder(json.JSONEncoder):
    def default(self, obj, **kwargs):
        if isinstance(obj, (ObjectId,)):
            return str(obj)
        return json.JSONEncoder.default(self, obj, **kwargs)


def get_user_collections(username, db):
    # Find collections for the specified username
    username_prefix = f'{username}_'
    collections_with_username = [collection for collection in db.list_collection_names() if collection.startswith(username_prefix)]
    return collections_with_username


def get_first_slide_from_collections(username, db):
    collections_with_username = get_user_collections(username, db)
    result_data = []

    for collection_name in collections_with_username:
        collection = db[collection_name]
        first_slide_document = collection.find_one({'slide_number': 1})

        if first_slide_document:
            # Ensure the '_id' field is converted to a string for JSON serialization
            first_slide_document['_id'] = str(first_slide_document['_id'])
            result_data.append(first_slide_document)

            topic = first_slide_document.get('topic', '')

    return result_data


def get_images_in_folder(topic):
    """
    Get a list of image files in the specified folder.

    Args:
        folder_path (str): The path to the folder containing image files.

    Returns:
        List[str]: A list of image file names.
    """
    allowed_extensions = {'.jpg', '.jpeg', '.png', '.gif'}  # Add more extensions if needed
    image_files = []

    try:
        for file_name in os.listdir('static'):
            file_path = os.path.join('static', {topic})

            # Check if it's a file and has an allowed extension
            if os.path.isfile(file_path) and os.path.splitext(file_name)[1].lower() in allowed_extensions:
                image_files.append(file_name)

    except Exception as e:
        print(f"Error getting images from folder: {e}")

    return image_files

def view_images(folder_name):
    folder_path = os.path.join('static', folder_name)
    image_files = get_images_in_folder(folder_path)
    return render_template('image_gallery.html', images=image_files)

@app.route('/search_collections')
def search_collections():
    if 'username' in session:
        username = session['username']
        result_data = get_first_slide_from_collections(username, db)

        # Print success if the search is done
        print('success')

        # Convert Base64 image data
        for document in result_data:
            if 'image_path' in document:
                try:
                    # Read the image file and encode it in Base64
                    image_path = document['image_path']
                    with open(image_path, 'rb') as image_file:
                        image_data = base64.b64encode(image_file.read()).decode('utf-8')
                        document['image_data'] = f"data:image/jpeg;base64,{image_data}"
                except Exception as e:
                    print(f"Error encoding image data: {e}")

        return render_template('Home.html', data=result_data)
    else:
        return jsonify({'error': 'User not logged in'})

@app.route('/get_image/<document_id>')
def get_image(document_id):
    # Retrieve the document by its ID
    document = collection.find_one({"_id": document_id})

    if document:
        # Extract the image data from the document
        image_data = document.get("image_data", None)

        if image_data:
            # Convert the image data to bytes
            image_bytes = io.BytesIO(image_data)

            # Return the image as a response with the correct content type
            return send_file(image_bytes, mimetype='image/jpeg')

    # Return an error response if the document or image data is not found
    return "Image not found", 404


@app.route('/view_presentation/<folder_path>')
def view_presentation(folder_path):
    try:
        # Get the folder_path from the query parameters
        folder_path = request.args.get('folder_path')

        if not folder_path:
            raise ValueError("Folder path not provided")

        # Construct the full path to the folder based on the provided folder_path
        full_folder_path = os.path.join('static', folder_path)

        # Get the list of image files in the specified folder_path
        image_files = get_image_files(full_folder_path)

        return render_template('ViewPresentation.html', image_files=image_files)

    except Exception as e:
        print(f"Error viewing presentation: {e}")
        return "Presentation not found", 404


def get_image_files(folder_path):
    try:
        # Use os.listdir to get a list of files in the specified folder_path
        image_files = [f for f in os.listdir(folder_path) if f.lower().endswith(('.png', '.jpg', '.jpeg'))]
        # Construct the full path for each image file
        image_files = [os.path.join(folder_path, file) for file in image_files]
        return image_files
    except Exception as e:
        print(f"Error getting image files: {e}")
        return []





@app.route('/logout')
@login_required
def logout():
    logout_user()
    return redirect(url_for('signin'))


@app.route('/download/<filename>', methods=['GET'])
def download_file(filename):
    try:
        return send_from_directory('generated', filename, as_attachment=True)

    except FileNotFoundError:
        os.abort(404)




app.config['PERMANENT_SESSION_LIFETIME'] = 24 * 60 * 60
if __name__ == '__main__':
    app.run(debug=True)