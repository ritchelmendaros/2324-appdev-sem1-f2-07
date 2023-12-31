import math
import os
from utils.template_simple import simple, update_simple
from utils.template_dark_modern import dark_modern, update_dark_modern
from utils.template_minimal_blue import minimal_blue, update_minimal_blue
from utils.template_minimal_darkgreen import minimal_darkgreen, update_minimal_darkgreen
from utils.template_minimal_neon import minimal_neon, update_minimal_neon
from utils.template_minimal_gray import minimal_gray, update_minimal_gray
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

dir_path = 'static/presentations'


def parse_response(response):
    holdTitle = ""
    isOneNewLine = 0
    isSlideTitleNewLine = 0
    # select split structure
    if response.__contains__('\n\n\n'):  # Ends with 2 newlines
        slides = response.split('\n\n\n')
    else:  # End with 1 newline
        isOneNewLine = 1
        slides = response.split('\n\n')

    slides_content = []

    for slide in slides:
        lines = slide.split('\n')

        if isOneNewLine == 1 and (len(lines) == 1 or isSlideTitleNewLine == 1): # End with 1 newline it have newline after slide title
            if holdTitle == "":
                title_line = lines[0]
                title = title_line.split(': ', 1)[1] if ': ' in title_line else title_line
                isSlideTitleNewLine = 1
                holdTitle = title
                continue
            else:
                title = holdTitle
                holdTitle = ""
        else:
            title_line = lines[0]
            title = title_line.split(': ', 1)[1] if ': ' in title_line else title_line

        #Contents for slide that ends with 2 newlines
        if isOneNewLine == 0:
            if lines.__contains__('\n\n'):# have a newline after the slide title
                content_lines = [line.lstrip('- ') for line in lines[2:] if line and 'Content:' not in line]
            else:
                content_lines = [line.lstrip('- ') for line in lines[1:] if line and 'Content:' not in line]
        else: #Contents for slide that ends with 1 newline
            if isSlideTitleNewLine == 1: # There is a newline after slide title
                content_lines = [line.lstrip('- ') for line in lines[0:] if line and 'Content:' not in line]
            else:
                content_lines = [line.lstrip('- ') for line in lines[1:] if line and 'Content:' not in line]

            isSlideTitleNewLine = 0

        content = '\n'.join(content_lines)
        keyword_line = [line for line in lines if 'Keyword:' in line or 'Keywords:' in line]
        keyword = keyword_line[0].split(': ', 1)[1].strip() if keyword_line else 'computer'
        slides_content.append({'title': title, 'content': content, 'keyword': keyword})

    return slides_content


def delete_first_two_slides(presentation):
    slide_ids = [1, 0, len(presentation.slides)-1]
    for slide_id in slide_ids:
        if slide_id < len(presentation.slides):
            xml_slides = presentation.slides._sldIdLst
            slides = list(xml_slides)
            xml_slides.remove(slides[slide_id])


def create_ppt(slides_content, template_choice, presentation_title, presenter_name):
    template_path = os.path.join(dir_path, f"{template_choice}.pptx")

    prs = Presentation(template_path)
    title_slide_layout = prs.slide_layouts[0]

    # add title slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = presentation_title

    #add subtitle
    subtitle = slide.placeholders[1]
    subtitle.text = f"Presented by {presenter_name}"

    if template_choice == 'simple':
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(90)
                run.font.name = 'Gill Sans MT'
                run.font.color.rgb = RGBColor(5, 14, 56) # RGB for orange color
        for paragraph in subtitle.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(40)
                run.font.name = 'Gill Sans MT'
                run.font.color.rgb = RGBColor(5, 14, 56) # RGB for orange color
        simple(prs, slides_content)

    elif template_choice == 'dark_modern':
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Montserrat'
                run.font.size = Pt(90)
                run.font.color.rgb = RGBColor(114, 222, 173)  # RGB for orange color
        for paragraph in subtitle.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Montserrat'
                run.font.size = Pt(40)
                run.font.color.rgb = RGBColor(255, 255, 255)  # RGB for orange color
        dark_modern(prs, slides_content)

    elif template_choice == 'minimal_blue':
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Roboto Slab'
                run.font.size = Pt(90)
                run.font.color.rgb = RGBColor(255, 255, 255)
        for paragraph in subtitle.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Roboto Slab'
                run.font.size = Pt(40)
                run.font.color.rgb = RGBColor(139, 195, 74)
        minimal_blue(prs, slides_content)

    elif template_choice == 'minimal_darkgreen':
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Oswald'
                run.font.size = Pt(90)
                run.font.color.rgb = RGBColor(255, 255, 255)
        for paragraph in subtitle.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(40)
                run.font.name = 'Oswald'
                run.font.color.rgb = RGBColor(255, 255, 255) # RGB for orange color
        minimal_darkgreen(prs, slides_content)

    elif template_choice == 'minimal_neon':
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Gill Sans MT'
                run.font.size = Pt(90)
                run.font.color.rgb = RGBColor(0, 0, 0)
        for paragraph in subtitle.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(40)
                run.font.name = 'Gill Sans MT'
                run.font.color.rgb = RGBColor(0, 0, 0) # RGB for orange color
        minimal_neon(prs, slides_content)

    elif template_choice == 'minimal_gray':
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Britannic Bold'
                run.font.size = Pt(90)
                run.font.color.rgb = RGBColor(0, 0, 0)
        for paragraph in subtitle.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(40)
                run.font.name = 'Franklin Gothic Demi Cond'
                run.font.color.rgb = RGBColor(0, 0, 0) # RGB for orange color
        minimal_gray(prs, slides_content)

    # Delete the first two slides after all new slides have been added
    delete_first_two_slides(prs)

    # Save the presentation
    prs.save(os.path.join('generated', 'generated_presentation.pptx'))


def update_slide_ppt(slides_content, file_path, auto, hasPicture, template_choice, slideNum, retain):
    ppt = os.path.join('generated', 'generated_presentation.pptx')
    prs = Presentation(ppt)
    try:
        num = int(slideNum)
        print(num)

        # add content slides
        if num < len(prs.slides):
            remove_all_elements(prs, num, retain, auto)
        else:
            prs.slides.add_slide(prs.slide_layouts[6])

        if template_choice == 'simple':
            update_simple(prs, file_path, auto, hasPicture, slides_content[0],num)
        elif template_choice == 'dark_modern':
            update_dark_modern(prs, file_path, auto, hasPicture, slides_content[0], num)
        elif template_choice == 'minimal_blue':
            update_minimal_blue(prs, file_path, auto, hasPicture, slides_content[0], num)
        elif template_choice == 'minimal_darkgreen':
            update_minimal_darkgreen(prs, file_path, auto, hasPicture, slides_content[0], num)
        elif template_choice == 'minimal_neon':
            update_minimal_neon(prs, file_path, auto, hasPicture, slides_content[0], num)
        elif template_choice == 'minimal_gray':
            update_minimal_gray(prs, file_path, auto, hasPicture, slides_content[0], num)

        prs.save(os.path.join('generated', 'generated_presentation.pptx'))
    finally:
        print("yees")


def remove_all_elements(prs, slideNum, retain, auto):
    # Iterate through each shape on the slide and remove it
    slide = prs.slides[slideNum]
    i = 0
    for shape in slide.shapes:
        if i == 0 and retain and not auto:
            retain = False
            continue
        else:
            shape.element.getparent().remove(shape.element)

    slide.notes_slide.notes_text_frame.clear()