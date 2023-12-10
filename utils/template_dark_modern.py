import math
import os
from io import BytesIO
import requests
from pptx.util import Inches
from utils.tools import slide_format,search_pexels_images


def dark_modern(prs, slides_content):
    numslides = len(slides_content)
    divided = (numslides - 1) / 4

    first = math.ceil(divided + 1)
    second = math.ceil(first + divided)
    third = math.ceil(second + divided)

    count = 1
    for slide_content in slides_content:
        if count < first:
            slide = prs.slides.add_slide(prs.slide_layouts[10])
            # file_path = os.path.join('static', 'pictures', '4.png')
            image_url = search_pexels_images(slide_content['title'])
            if image_url:
                # Download the image
                response = requests.get(image_url)
                image_data = BytesIO(response.content)
                slide.shapes.add_picture(image_data, Inches(2.61), Inches(1.00), Inches(7.90), Inches(4.52))
            title_box = slide.shapes.add_textbox(Inches(2.61), Inches(5.95), Inches(7.81), Inches(4.1))
            content_box = slide.shapes.add_textbox(Inches(10.79), Inches(1.00), Inches(8.92), Inches(8.7))
            title_box.text = slide_content['title']
            content_box.text = slide_content['content']

            title_box.text_frame.word_wrap = content_box.text_frame.word_wrap = True
            content_frame = content_box.text_frame
            title_frame = title_box.text_frame

            slide_format(title_frame, 50, 'Montserrat', 114, 222, 173, 0, 0)
            slide_format(content_frame, 32, 'Lato', 255, 255, 255, 0, 16)

        elif count < second and count >= first:
            slide = prs.slides.add_slide(prs.slide_layouts[10])
            # file_path = os.path.join('static', 'pictures', '4.png')
            image_url = search_pexels_images(slide_content['title'])
            if image_url:
                # Download the image
                response = requests.get(image_url)
                image_data = BytesIO(response.content)
                slide.shapes.add_picture(image_data, Inches(0.79), Inches(3.43), Inches(7.15), Inches(7.15))
            title_box = slide.shapes.add_textbox(Inches(2.40), Inches(1.04), Inches(17.22), Inches(1.31))
            content_box = slide.shapes.add_textbox(Inches(9.62), Inches(2.86), Inches(9.78), Inches(7.1))
            title_box.text = slide_content['title']
            content_box.text = slide_content['content']

            title_box.text_frame.word_wrap = content_box.text_frame.word_wrap = True
            content_frame = content_box.text_frame
            title_frame = title_box.text_frame

            slide_format(title_frame, 50, 'Montserrat', 114, 222, 173, 0, 0)
            slide_format(content_frame, 32, 'Lato', 255, 255, 255, 0, 20)

        elif count < third and count >= second:
            slide = prs.slides.add_slide(prs.slide_layouts[10])
            title_box = slide.shapes.add_textbox(Inches(2.40), Inches(1.04), Inches(18), Inches(2))
            content_box = slide.shapes.add_textbox(Inches(1.54), Inches(3.43), Inches(18), Inches(7.1))
            title_box.text = slide_content['title']
            content_box.text = slide_content['content']

            title_box.text_frame.word_wrap = content_box.text_frame.word_wrap = True
            content_frame = content_box.text_frame
            title_frame = title_box.text_frame

            slide_format(title_frame, 50, 'Montserrat', 114, 222, 173, 0, 0)
            slide_format(content_frame, 32, 'Lato', 255, 255, 255, 0, 25)

        else:
            # image_url = search_pexels_images(slide_content['title'])
            slide = prs.slides.add_slide(prs.slide_layouts[10])
            # file_path = os.path.join('static', 'pictures', '4.png')
            image_url = search_pexels_images(slide_content['title'])
            if image_url:
                # Download the image
                response = requests.get(image_url)
                image_data = BytesIO(response.content)
                slide.shapes.add_picture(image_data, Inches(11.60), Inches(2.50), Inches(7.05), Inches(7.81))
            title_box = slide.shapes.add_textbox(Inches(2.40), Inches(1.04), Inches(9.71), Inches(2.12))
            content_box = slide.shapes.add_textbox(Inches(0.75), Inches(3.43), Inches(9.71), Inches(7))
            title_box.text = slide_content['title']
            content_box.text = slide_content['content']

            title_box.text_frame.word_wrap = content_box.text_frame.word_wrap = True
            content_frame = content_box.text_frame
            title_frame = title_box.text_frame

            slide_format(title_frame, 66, 'Montserrat', 114, 222, 173, 0, 0)
            slide_format(content_frame, 32, 'Lato', 255, 255, 255, 0, 20)

        count+=1


def update_dark_modern(prs, file_path, auto, hasPicture, slide_content, slideNum):
    numslides = len(prs.slides)
    divided = (numslides - 1) / 4

    first = math.ceil(divided + 1)
    second = math.ceil(first + divided)
    third = math.ceil(second + divided)

    if slideNum < first:
        slide = prs.slides[slideNum]
        if auto:
            image_url = search_pexels_images(slide_content['title'])
            if image_url:
                # Download the image
                response = requests.get(image_url)
                image_data = BytesIO(response.content)
                slide.shapes.add_picture(image_data, Inches(2.61), Inches(1.00), Inches(7.90), Inches(4.52))
        elif hasPicture:
            slide.shapes.add_picture(file_path, Inches(2.61), Inches(1.00), Inches(7.90), Inches(4.52))

        title_box = slide.shapes.add_textbox(Inches(2.61), Inches(5.95), Inches(7.81), Inches(4.1))
        title_box.text = slide_content['title']
        content_box = slide.shapes.add_textbox(Inches(10.79), Inches(1.00), Inches(8.92), Inches(8.7))
        content_box.text = slide_content['content']

        title_box.text_frame.word_wrap = content_box.text_frame.word_wrap = True
        content_frame = content_box.text_frame
        title_frame = title_box.text_frame

        slide_format(title_frame, 50, 'Montserrat', 114, 222, 173, 0, 0)
        slide_format(content_frame, 32, 'Lato', 255, 255, 255, 0, 16)

    elif slideNum < second and slideNum >= first:
        slide = prs.slides[slideNum]
        if auto:
            image_url = search_pexels_images(slide_content['title'])
            if image_url:
                # Download the image
                response = requests.get(image_url)
                image_data = BytesIO(response.content)
                slide.shapes.add_picture(image_data, Inches(0.79), Inches(3.43), Inches(7.15), Inches(7.15))
        elif hasPicture:
            slide.shapes.add_picture(file_path, Inches(0.79), Inches(3.43), Inches(7.15), Inches(7.15))

        title_box = slide.shapes.add_textbox(Inches(2.40), Inches(1.04), Inches(17.22), Inches(1.31))
        content_box = slide.shapes.add_textbox(Inches(9.62), Inches(2.86), Inches(9.78), Inches(7.1))
        title_box.text = slide_content['title']
        content_box.text = slide_content['content']

        title_box.text_frame.word_wrap = content_box.text_frame.word_wrap = True
        content_frame = content_box.text_frame
        title_frame = title_box.text_frame

        slide_format(title_frame, 50, 'Montserrat', 114, 222, 173, 0, 0)
        slide_format(content_frame, 32, 'Lato', 255, 255, 255, 0, 20)

    elif slideNum < third and slideNum >= second:
        slide = prs.slides[slideNum]
        title_box = slide.shapes.add_textbox(Inches(2.40), Inches(1.04), Inches(18), Inches(2))
        content_box = slide.shapes.add_textbox(Inches(1.54), Inches(3.43), Inches(18), Inches(7.1))
        title_box.text = slide_content['title']
        content_box.text = slide_content['content']

        title_box.text_frame.word_wrap = content_box.text_frame.word_wrap = True
        content_frame = content_box.text_frame
        title_frame = title_box.text_frame

        slide_format(title_frame, 50, 'Montserrat', 114, 222, 173, 0, 0)
        slide_format(content_frame, 32, 'Lato', 255, 255, 255, 0, 25)

    else:
        slide = prs.slides[slideNum]
        if auto:
            image_url = search_pexels_images(slide_content['title'])
            if image_url:
                # Download the image
                response = requests.get(image_url)
                image_data = BytesIO(response.content)
                slide.shapes.add_picture(image_data, Inches(11.60), Inches(2.50), Inches(7.05), Inches(7.81))
        elif hasPicture:
            slide.shapes.add_picture(file_path, Inches(11.60), Inches(2.50), Inches(7.05), Inches(7.81))

        title_box = slide.shapes.add_textbox(Inches(2.40), Inches(1.04), Inches(9.71), Inches(2.12))
        content_box = slide.shapes.add_textbox(Inches(0.75), Inches(3.43), Inches(9.71), Inches(7))
        title_box.text = slide_content['title']
        content_box.text = slide_content['content']

        title_box.text_frame.word_wrap = content_box.text_frame.word_wrap = True
        content_frame = content_box.text_frame
        title_frame = title_box.text_frame

        slide_format(title_frame, 50, 'Montserrat', 114, 222, 173, 0, 0)
        slide_format(content_frame, 32, 'Lato', 255, 255, 255, 0, 20)